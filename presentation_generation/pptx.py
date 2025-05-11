import logging
import os
import torch

from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt

log = logging.getLogger(__name__)


class Pptx:

    duration_template = "oxml_duration_template.xml"
    pixels_per_inch = 96  # Subjective fixed value

    def __init__(self, output_folder):
        self.output_folder = output_folder
        Path(self.output_folder).mkdir(parents=True, exist_ok=True)

    def get_next_filepath(self, title):
        """
        Generate the next filepath available by appending `_X`
        """
        extension = 'pptx'
        filename = title.lower().replace(' ', '_')
        filepath = f"{self.output_folder}/{filename}.{extension}"
        if os.path.isfile(filepath):
            log.info(f"file exsists {filepath}")
            i = 1
            filepath = f"{self.output_folder}/{filename}_{i}.{extension}"
            while os.path.isfile(filepath):
                log.info(filepath)
                i = i+1
                filepath = f"{self.output_folder}/{filename}_{i}.{extension}"
        return filepath

    def get_duration_fragment(self, duration=20):
        """
        Get the duration xml fragment to add to each slide with automatic transition
        duration: slide duration in seconds
        returns: the xml fragment to add to the slide
        """
        template_path = os.path.join(os.path.dirname(__file__), self.duration_template)
        with open(template_path) as f:
            xml_fragment = f.read()
        xml_fragment = xml_fragment.replace("{duration}", str(duration * 1000))
        return parse_xml(xml_fragment)

    def add_presentation_title(self, root, title):
        """
        Generate the first presentaiton slide
        """
        title_slide = root.slides.add_slide(root.slide_layouts[0])

        main_title = title_slide.shapes.title
        main_title.width = root.slide_width
        main_title.left = 0
        main_title.top = root.slide_height//2
        main_title.text = title
        main_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(188, 194, 200)

        sub_title = title_slide.placeholders[1]
        sub_title.width = root.slide_width
        sub_title.left = 0
        sub_title.top = root.slide_height - Pt(12) - Inches(0.2)
        sub_title.text = "Generated using ComfyUI-PechaKucha"
        sub_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(188, 194, 200)
        sub_title.text_frame.paragraphs[0].font.size = Pt(12)

        background = title_slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(2, 3, 3)

    def generate(self, title, slide_duration, images, show_titles, titles):
        """
        Generate the powerpoint presentation file
        Ref for slide types:
        0 ->  title and subtitle
        1 ->  title and content
        2 ->  section header
        3 ->  two content
        4 ->  Comparison
        5 ->  Title only
        6 ->  Blank
        7 ->  Content with caption
        8 ->  Pic with caption
        """
        log.debug(f"{title} Pecha Kucha generation with {len(images)} images and duration {slide_duration}")

        root = Presentation()
        first_image = Image.open(images[0])
        width, height = first_image.size
        root.slide_width = Inches(width//self.pixels_per_inch)
        root.slide_height = Inches(height//self.pixels_per_inch)

        self.add_presentation_title(root, title)

        if show_titles:
            assert(len(images) == len(titles))

        slide_layout = root.slide_layouts[0]
        left = top = Inches(0)
        width = root.slide_width

        for i, image in enumerate(images):
            slide = root.slides.add_slide(slide_layout)

            pic = slide.shapes.add_picture(image, left, top, width)
            # Send picture to background
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)

            if show_titles:
                slide.shapes.title.text = titles[i]
                slide.shapes.title.width = root.slide_width
                slide.shapes.title.left = 0
                slide.shapes.title.top = root.slide_height - Pt(16) - Inches(0.5)
                slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(2, 3, 3)
                slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(16)

            # Adding automatic transition
            slide.element.insert(-1, self.get_duration_fragment(slide_duration))

        filepath = self.get_next_filepath(title)
        root.save(filepath)

        return filepath
